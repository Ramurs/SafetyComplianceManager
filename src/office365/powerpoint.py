from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.config import get_settings


def generate_executive_summary(audits: list, risks: list) -> Path:
    """Generate a PowerPoint executive summary presentation."""
    settings = get_settings()
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Compliance & Risk Executive Summary"
    slide.placeholders[1].text = "Safety Compliance Manager"

    # Overview slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Overview"
    body = slide.placeholders[1]
    tf = body.text_frame

    total_audits = len(audits)
    completed_audits = sum(1 for a in audits if a.status == "completed")
    total_risks = len(risks)
    high_risks = sum(1 for r in risks if r.score >= 16)
    medium_risks = sum(1 for r in risks if 9 <= r.score < 16)

    tf.text = f"Total Audits: {total_audits} ({completed_audits} completed)"
    tf.add_paragraph().text = f"Total Risks: {total_risks}"
    tf.add_paragraph().text = f"High Risks (16+): {high_risks}"
    tf.add_paragraph().text = f"Medium Risks (9-15): {medium_risks}"

    # Top Risks slide
    if risks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Top Risks"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for risk in sorted(risks, key=lambda r: r.score, reverse=True)[:5]:
            p = tf.add_paragraph()
            p.text = f"[Score: {risk.score}] {risk.title}"
            p.font.size = Pt(14)
            if risk.score >= 16:
                p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
            elif risk.score >= 9:
                p.font.color.rgb = RGBColor(0xFF, 0x99, 0x00)

    # Recent Audits slide
    if audits:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Recent Audits"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for audit in audits[:5]:
            p = tf.add_paragraph()
            p.text = f"[{audit.status.upper()}] {audit.title}"
            p.font.size = Pt(14)

    output_path = settings.output_dir / "executive_summary.pptx"
    prs.save(str(output_path))
    return output_path
